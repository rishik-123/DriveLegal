import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    BG_COLOR = RGBColor(11, 15, 26)         # Deep slate/navy dark (#0B0F1A)
    CARD_BG_COLOR = RGBColor(21, 28, 45)    # Slate card background (#151C2D)
    TEXT_MAIN = RGBColor(248, 250, 252)     # Off-white (#F8FAFC)
    TEXT_MUTED = RGBColor(148, 163, 184)    # Slate grey (#94A3B8)
    ACCENT_CYAN = RGBColor(6, 182, 212)     # Electric Cyan (#06B6D4)
    ACCENT_INDIGO = RGBColor(99, 102, 241)  # Vivid Indigo (#6366F1)
    ACCENT_GREEN = RGBColor(16, 185, 129)   # Emerald (#10B981)
    ACCENT_RED = RGBColor(239, 68, 68)      # Coral Red (#EF4444)
    
    # Fonts
    FONT_TITLE = 'Trebuchet MS'
    FONT_BODY = 'Calibri'
    
    # Image Paths
    LANDING_IMG = 'landing_page.png'
    LOGIN_IMG = 'login_page.png'
    
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    def add_header(slide, title_text, category_text="DRIVEVERSE"):
        # Add thin top accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = ACCENT_CYAN
        accent_bar.line.fill.background()
        
        # Category label (small, uppercase above title)
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = FONT_BODY
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_INDIGO
        
        # Main title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_MAIN

    # ==========================================
    # SLIDE 1: INTRODUCTION / WELCOME
    # ==========================================
    slide_layout = prs.slide_layouts[6] # blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide1)
    
    # Central Card layout
    center_card = slide1.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5)
    )
    center_card.fill.solid()
    center_card.fill.fore_color.rgb = CARD_BG_COLOR
    center_card.line.color.rgb = ACCENT_INDIGO
    center_card.line.width = Pt(2)
    
    # Title & Subtitle
    title_box = slide1.shapes.add_textbox(Inches(2.0), Inches(2.2), Inches(9.333), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "🚗 DriveVerse"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MAIN
    
    sub_box = slide1.shapes.add_textbox(Inches(2.0), Inches(3.4), Inches(9.333), Inches(0.8))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p2 = tf_sub.paragraphs[0]
    p2.text = "One Identity. Every Vehicle Service."
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = FONT_BODY
    p2.font.size = Pt(20)
    p2.font.italic = True
    p2.font.color.rgb = ACCENT_CYAN
    
    # Team Info
    team_box = slide1.shapes.add_textbox(Inches(2.0), Inches(4.4), Inches(9.333), Inches(1.0))
    tf_team = team_box.text_frame
    tf_team.word_wrap = True
    p3 = tf_team.paragraphs[0]
    p3.text = "Presented by:"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = FONT_BODY
    p3.font.size = Pt(13)
    p3.font.color.rgb = TEXT_MUTED
    
    p4 = tf_team.add_paragraph()
    p4.text = "Rishik Jariwala  &  Jash Chothani"
    p4.alignment = PP_ALIGN.CENTER
    p4.font.name = FONT_TITLE
    p4.font.size = Pt(18)
    p4.font.bold = True
    p4.font.color.rgb = TEXT_MAIN

    # ==========================================
    # SLIDE 2: PROBLEM STATEMENT & EXISTING GAPS
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "The Problem & Existing Gaps", "01. PROBLEM & MARKET LANDSCAPE")
    
    col_width = 5.6
    col_gap = 0.533
    
    # Left Column: Problems
    card_l = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(col_width), Inches(4.7)
    )
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = CARD_BG_COLOR
    card_l.line.color.rgb = ACCENT_RED
    card_l.line.width = Pt(1.5)
    
    tb_l = slide2.shapes.add_textbox(Inches(1.1), Inches(2.05), Inches(col_width - 0.6), Inches(4.2))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = 0
    
    p_l_t = tf_l.paragraphs[0]
    p_l_t.text = "Problem Statement"
    p_l_t.font.name = FONT_TITLE
    p_l_t.font.size = Pt(20)
    p_l_t.font.bold = True
    p_l_t.font.color.rgb = ACCENT_CYAN
    p_l_t.space_after = Pt(12)
    
    problems = [
      ("Information Silos", "Documents (DL, RC, Insurance, PUC) are scattered across physical files, email threads, and disjointed portals."),
      ("Compliance Failures", "Missed renewal deadlines result in heavy traffic fines and legal disputes."),
      ("Security Risks", "Lack of geofencing triggers or real-time location alerts leaves vehicles vulnerable to theft.")
    ]
    for p_title, p_desc in problems:
        p_pt = tf_l.add_paragraph()
        p_pt.space_before = Pt(8)
        p_pt.text = "⚠️ " + p_title
        p_pt.font.name = FONT_TITLE
        p_pt.font.size = Pt(14)
        p_pt.font.bold = True
        p_pt.font.color.rgb = TEXT_MAIN
        
        p_pd = tf_l.add_paragraph()
        p_pd.text = p_desc
        p_pd.font.name = FONT_BODY
        p_pd.font.size = Pt(12)
        p_pd.font.color.rgb = TEXT_MUTED

    # Right Column: Existing Gaps
    card_r = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + col_width + col_gap), Inches(1.8), Inches(col_width), Inches(4.7)
    )
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = CARD_BG_COLOR
    card_r.line.color.rgb = ACCENT_RED
    card_r.line.width = Pt(1.5)
    
    tb_r = slide2.shapes.add_textbox(Inches(0.8 + col_width + col_gap + 0.3), Inches(2.05), Inches(col_width - 0.6), Inches(4.2))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = tf_r.margin_bottom = 0
    
    p_r_t = tf_r.paragraphs[0]
    p_r_t.text = "Gaps in Existing Solutions"
    p_r_t.font.name = FONT_TITLE
    p_r_t.font.size = Pt(20)
    p_r_t.font.bold = True
    p_r_t.font.color.rgb = ACCENT_RED
    p_r_t.space_after = Pt(12)
    
    gaps = [
      ("Government Portals (DigiLocker)", "Static databases with zero proactive alert triggers. Poor user experience and no navigational utilities."),
      ("Standalone Apps", "Cluttered with advertisements. Single-purpose and fail to link directly with government records."),
      ("Physical glove-box files", "High risk of damage or loss; completely ineffective in security alerts or on-road emergencies.")
    ]
    for g_title, g_desc in gaps:
        p_gt = tf_r.add_paragraph()
        p_gt.space_before = Pt(8)
        p_gt.text = "✗ " + g_title
        p_gt.font.name = FONT_TITLE
        p_gt.font.size = Pt(14)
        p_gt.font.bold = True
        p_gt.font.color.rgb = TEXT_MAIN
        
        p_gd = tf_r.add_paragraph()
        p_gd.text = g_desc
        p_gd.font.name = FONT_BODY
        p_gd.font.size = Pt(12)
        p_gd.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 3: TECH STACK & USER FLOW (with Screenshot)
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "Tech Stack & User Flow", "02. ARCHITECTURE & TECHNOLOGY")
    
    # Left column: Tech Stack details
    tech_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_tech = tech_box.text_frame
    tf_tech.word_wrap = True
    
    p_tech_h = tf_tech.paragraphs[0]
    p_tech_h.text = "The Technology Stack"
    p_tech_h.font.name = FONT_TITLE
    p_tech_h.font.size = Pt(20)
    p_tech_h.font.bold = True
    p_tech_h.font.color.rgb = ACCENT_CYAN
    p_tech_h.space_after = Pt(14)
    
    tech_items = [
      ("Frontend", "Next.js 15 · TypeScript · TailwindCSS · Framer Motion"),
      ("Backend", "FastAPI · Python 3.12 · SQLAlchemy (async) · PostgreSQL · Redis"),
      ("AI Copilot", "Google Gemini 2.0 Flash (Astra AI integration)"),
      ("Security", "JWT Token Rotation · bcrypt password hashing · AES Fernet encryption")
    ]
    for t_name, t_detail in tech_items:
        p_ti = tf_tech.add_paragraph()
        p_ti.space_before = Pt(8)
        p_ti.text = f"💻 {t_name}: "
        p_ti.font.name = FONT_TITLE
        p_ti.font.size = Pt(13)
        p_ti.font.bold = True
        p_ti.font.color.rgb = TEXT_MAIN
        
        p_td = tf_tech.add_paragraph()
        p_td.text = f"    {t_detail}"
        p_td.font.name = FONT_BODY
        p_td.font.size = Pt(12)
        p_td.font.color.rgb = TEXT_MUTED

    # Right column: User Flow & Screenshot
    flow_box = slide3.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.0))
    tf_flow = flow_box.text_frame
    tf_flow.word_wrap = True
    
    p_flow_h = tf_flow.paragraphs[0]
    p_flow_h.text = "User & Data Flow"
    p_flow_h.font.name = FONT_TITLE
    p_flow_h.font.size = Pt(20)
    p_flow_h.font.bold = True
    p_flow_h.font.color.rgb = ACCENT_INDIGO
    
    p_flow_b = tf_flow.add_paragraph()
    p_flow_b.space_before = Pt(6)
    p_flow_b.text = "1. User logins via OTP/OAuth ➔ 2. Garage pulls live DigiLocker & e-Challan data ➔ 3. Astra AI assists user and displays telemetry on maps."
    p_flow_b.font.name = FONT_BODY
    p_flow_b.font.size = Pt(12)
    p_flow_b.font.color.rgb = TEXT_MUTED
    
    # Add Landing Page Screenshot
    if os.path.exists(LANDING_IMG):
        # Image Frame Border
        img_border = slide3.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(3.7), Inches(5.7), Inches(2.8)
        )
        img_border.fill.solid()
        img_border.fill.fore_color.rgb = CARD_BG_COLOR
        img_border.line.color.rgb = ACCENT_INDIGO
        img_border.line.width = Pt(1.5)
        
        # Add actual image inside border
        slide3.shapes.add_picture(
            LANDING_IMG, Inches(6.82), Inches(3.72), Inches(5.66), Inches(2.76)
        )

    # ==========================================
    # SLIDE 4: OUR SOLUTION (with Login Screenshot)
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "Our Solution: DriveVerse", "03. THE UNIFIED ECOSYSTEM")
    
    # Left column: Core Pillars
    sol_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_sol = sol_box.text_frame
    tf_sol.word_wrap = True
    
    p_sol_h = tf_sol.paragraphs[0]
    p_sol_h.text = "Core Pillars of DriveVerse"
    p_sol_h.font.name = FONT_TITLE
    p_sol_h.font.size = Pt(20)
    p_sol_h.font.bold = True
    p_sol_h.font.color.rgb = ACCENT_CYAN
    p_sol_h.space_after = Pt(14)
    
    pillars = [
      ("Unified Garage", "Track all vehicles in one hub with 7 switchable custom themes."),
      ("Astra AI Co-Pilot", "Gemini-powered voice and multilingual (EN/HI/MR) assistant."),
      ("Integrated Doc Vault", "Secure DigiLocker vault tracking DL, RC, PUC, Insurance expiries."),
      ("Advanced Anti-Theft", "Geofencing safe zones, tamper detection, and real-time OTP alerts.")
    ]
    for p_title, p_desc in pillars:
        p_pt = tf_sol.add_paragraph()
        p_pt.space_before = Pt(8)
        p_pt.text = "⚡ " + p_title
        p_pt.font.name = FONT_TITLE
        p_pt.font.size = Pt(14)
        p_pt.font.bold = True
        p_pt.font.color.rgb = TEXT_MAIN
        
        p_pd = tf_sol.add_paragraph()
        p_pd.text = p_desc
        p_pd.font.name = FONT_BODY
        p_pd.font.size = Pt(12)
        p_pd.font.color.rgb = TEXT_MUTED

    # Right column: Authentication/Login UI Screenshot
    if os.path.exists(LOGIN_IMG):
        # Image Frame Border
        img_border = slide4.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5)
        )
        img_border.fill.solid()
        img_border.fill.fore_color.rgb = CARD_BG_COLOR
        img_border.line.color.rgb = ACCENT_INDIGO
        img_border.line.width = Pt(1.5)
        
        # Add actual image inside border
        slide4.shapes.add_picture(
            LOGIN_IMG, Inches(6.82), Inches(1.82), Inches(5.66), Inches(4.46)
        )

    # ==========================================
    # SLIDE 5: HOW IT HELPS TO AVOID THE PROBLEM
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "How DriveVerse Avoids the Problems", "04. SYSTEM WORKFLOW & BENEFITS")
    
    # Left Side: Bulleted features of prevention
    left_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    benefits = [
        ("Automated Alerts (No Expiry Oversights)", "Proactive 30/15/7 day expiry countdown alerts for Insurance & PUC policies via Email & SMS OTP ensure 100% compliance."),
        ("Geofencing Safe Zones (Active Theft Prevention)", "Real-time geofence circles alert vehicle owners instantly if their vehicle moves outside safe boundaries or detects tamper events."),
        ("Live Navigation (Avoid Speed Traps & Fines)", "Integrates real-time speed limit indicators, road safety warnings, and live traffic alerts, keeping the driver informed and ticket-free."),
        ("Unified Challan Center (No Delinquency)", "Direct e-Challan synchronization scans and alerts owners of traffic violations instantly, allowing immediate online resolution.")
    ]
    
    for idx, (title, desc) in enumerate(benefits):
        p_t = tf_left.add_paragraph() if idx > 0 else tf_left.paragraphs[0]
        if idx > 0:
            p_t.space_before = Pt(14)
        p_t.text = "🛡️ " + title
        p_t.font.name = FONT_TITLE
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_CYAN
        
        p_d = tf_left.add_paragraph()
        p_d.text = desc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_MUTED
        
    # Right Side: System architecture / data flow representation
    right_box = slide5.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_r_t = tf_right.paragraphs[0]
    p_r_t.text = "The DriveVerse Core Flow"
    p_r_t.font.name = FONT_TITLE
    p_r_t.font.size = Pt(20)
    p_r_t.font.bold = True
    p_r_t.font.color.rgb = TEXT_MAIN
    p_r_t.space_after = Pt(14)
    
    flows = [
        ("1. Data Sync", "FastAPI connects with DigiLocker, e-Challan APIs and Google Maps to collect active telemetry & compliance details."),
        ("2. Astra AI Intelligence", "Gemini AI parses user questions, navigation cues, and document status to formulate proactive driving recommendations."),
        ("3. Notification Hub", "Redis queues route critical notifications (expiries, tamper alerts, geofence breaches) directly to user devices.")
    ]
    
    for f_title, f_desc in flows:
        p_f_t = tf_right.add_paragraph()
        p_f_t.space_before = Pt(10)
        p_f_t.text = f_title
        p_f_t.font.name = FONT_TITLE
        p_f_t.font.size = Pt(14)
        p_f_t.font.bold = True
        p_f_t.font.color.rgb = ACCENT_INDIGO
        
        p_f_d = tf_right.add_paragraph()
        p_f_d.text = f_desc
        p_f_d.font.name = FONT_BODY
        p_f_d.font.size = Pt(12)
        p_f_d.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 6: SHOWSTOPPERS, CHALLENGES & FUTURE SCOPE
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide6)
    add_header(slide6, "Challenges & Future Vision", "05. ROADMAP & HURDLES")
    
    col_width_6 = 5.6
    col_gap_6 = 0.533
    
    # Left column: Challenges & Showstoppers
    card_l = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(col_width_6), Inches(4.7)
    )
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = CARD_BG_COLOR
    card_l.line.color.rgb = ACCENT_RED
    card_l.line.width = Pt(1.5)
    
    tb_l = slide6.shapes.add_textbox(Inches(1.1), Inches(2.05), Inches(col_width_6 - 0.6), Inches(4.2))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = 0
    
    p_l_t = tf_l.paragraphs[0]
    p_l_t.text = "Showstoppers & Challenges"
    p_l_t.font.name = FONT_TITLE
    p_l_t.font.size = Pt(20)
    p_l_t.font.bold = True
    p_l_t.font.color.rgb = ACCENT_CYAN
    p_l_t.space_after = Pt(14)
    
    challs = [
        ("Legacy Government API Access", "Integrating with regional RTO databases and state-specific e-Challan systems is complex due to slow response times and security protocols."),
        ("Offline Voice AI & Navigation", "Ensuring critical navigation alerts and Astra AI voice commands remain functional in areas with poor or zero cellular network connectivity.")
    ]
    for c_title, c_desc in challs:
        p_c_t = tf_l.add_paragraph()
        p_c_t.space_before = Pt(10)
        p_c_t.text = "⚡ " + c_title
        p_c_t.font.name = FONT_TITLE
        p_c_t.font.size = Pt(14)
        p_c_t.font.bold = True
        p_c_t.font.color.rgb = TEXT_MAIN
        
        p_c_d = tf_l.add_paragraph()
        p_c_d.text = c_desc
        p_c_d.font.name = FONT_BODY
        p_c_d.font.size = Pt(12)
        p_c_d.font.color.rgb = TEXT_MUTED

    # Right column: Future Scope
    card_r = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + col_width_6 + col_gap_6), Inches(1.8), Inches(col_width_6), Inches(4.7)
    )
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = CARD_BG_COLOR
    card_r.line.color.rgb = ACCENT_GREEN
    card_r.line.width = Pt(1.5)
    
    tb_r = slide6.shapes.add_textbox(Inches(0.8 + col_width_6 + col_gap_6 + 0.3), Inches(2.05), Inches(col_width_6 - 0.6), Inches(4.2))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = tf_r.margin_bottom = 0
    
    p_r_t = tf_r.paragraphs[0]
    p_r_t.text = "Future Scope"
    p_r_t.font.name = FONT_TITLE
    p_r_t.font.size = Pt(20)
    p_r_t.font.bold = True
    p_r_t.font.color.rgb = ACCENT_GREEN
    p_r_t.space_after = Pt(14)
    
    scopes = [
        ("Predictive Vehicle Diagnostics", "Deploying machine learning models to analyze engine telematics and predict parts wear-and-tear before a breakdown occurs."),
        ("Smart Electric Vehicle (EV) Portal", "Providing real-time battery analytics, health scoring, and smart route optimization mapped specifically through charging grids."),
        ("Logistics & Enterprise Fleet Portal", "Scaling the architecture to support logistics operators, enabling them to track and manage compliance for hundreds of trucks in a single dashboard.")
    ]
    for s_title, s_desc in scopes:
        p_s_t = tf_r.add_paragraph()
        p_s_t.space_before = Pt(8)
        p_s_t.text = "🚀 " + s_title
        p_s_t.font.name = FONT_TITLE
        p_s_t.font.size = Pt(14)
        p_s_t.font.bold = True
        p_s_t.font.color.rgb = TEXT_MAIN
        
        p_s_d = tf_r.add_paragraph()
        p_s_d.text = s_desc
        p_s_d.font.name = FONT_BODY
        p_s_d.font.size = Pt(11)
        p_s_d.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 7: THANK YOU SLIDE
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide7)
    
    # Center Card
    center_card7 = slide7.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5)
    )
    center_card7.fill.solid()
    center_card7.fill.fore_color.rgb = CARD_BG_COLOR
    center_card7.line.color.rgb = ACCENT_CYAN
    center_card7.line.width = Pt(2)
    
    # Title
    title_box7 = slide7.shapes.add_textbox(Inches(2.0), Inches(2.0), Inches(9.333), Inches(1.2))
    tf7 = title_box7.text_frame
    tf7.word_wrap = True
    p7_1 = tf7.paragraphs[0]
    p7_1.text = "Thank You!"
    p7_1.alignment = PP_ALIGN.CENTER
    p7_1.font.name = FONT_TITLE
    p7_1.font.size = Pt(64)
    p7_1.font.bold = True
    p7_1.font.color.rgb = TEXT_MAIN
    
    # Subtitle
    sub_box7 = slide7.shapes.add_textbox(Inches(2.0), Inches(3.2), Inches(9.333), Inches(0.8))
    tf_sub7 = sub_box7.text_frame
    tf_sub7.word_wrap = True
    p7_2 = tf_sub7.paragraphs[0]
    p7_2.text = "Redefining Indian Mobility, Redesigning Driving Experiences."
    p7_2.alignment = PP_ALIGN.CENTER
    p7_2.font.name = FONT_BODY
    p7_2.font.size = Pt(18)
    p7_2.font.italic = True
    p7_2.font.color.rgb = ACCENT_CYAN
    
    # Team members acknowledgement
    team_box7 = slide7.shapes.add_textbox(Inches(2.0), Inches(4.3), Inches(9.333), Inches(1.2))
    tf_team7 = team_box7.text_frame
    tf_team7.word_wrap = True
    p7_3 = tf_team7.paragraphs[0]
    p7_3.text = "Project Team:"
    p7_3.alignment = PP_ALIGN.CENTER
    p7_3.font.name = FONT_BODY
    p7_3.font.size = Pt(13)
    p7_3.font.color.rgb = TEXT_MUTED
    
    p7_4 = tf_team7.add_paragraph()
    p7_4.text = "Rishik Jariwala  |  Jash Chothani"
    p7_4.alignment = PP_ALIGN.CENTER
    p7_4.font.name = FONT_TITLE
    p7_4.font.size = Pt(18)
    p7_4.font.bold = True
    p7_4.font.color.rgb = TEXT_MAIN
    
    p7_5 = tf_team7.add_paragraph()
    p7_5.text = "DriveVerse — Built with Next.js, FastAPI, and Google Gemini AI"
    p7_5.alignment = PP_ALIGN.CENTER
    p7_5.font.name = FONT_BODY
    p7_5.font.size = Pt(11)
    p7_5.font.color.rgb = ACCENT_INDIGO
    p7_5.space_before = Pt(8)

    # Save presentation
    prs.save('DriveVerse_Presentation.pptx')
    print("Presentation generated successfully!")

if __name__ == '__main__':
    create_presentation()
